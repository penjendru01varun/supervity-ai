from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import io
import os
import datetime

class ContentAgent:
    def __init__(self, output_dir: str = "data/outputs"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def generate_pdf_report(self, incident_data: dict, research_data: dict) -> str:
        """
        Generate a professional PDF report with reportlab.
        """
        file_path = os.path.join(self.output_dir, f"Incident_Report_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf")
        c = canvas.Canvas(file_path, pagesize=letter)
        width, height = letter

        # Title
        c.setFont("Helvetica-Bold", 18)
        c.drawCentredString(width/2, height - 50, "ROAD SAFETY INCIDENT ANALYSIS REPORT")
        c.setFont("Helvetica", 12)
        c.drawCentredString(width/2, height - 70, f"Report ID: RSI-{datetime.datetime.now().strftime('%Y')}-001")
        
        # Section 1: Incident Summary
        c.setFont("Helvetica-Bold", 14)
        c.drawString(50, height - 120, "📍 INCIDENT SUMMARY")
        c.setFont("Helvetica", 11)
        c.drawString(50, height - 140, f"Location: {incident_data.get('location', {}).get('address', 'Unknown')}")
        c.drawString(50, height - 155, f"Date: {incident_data.get('datetime', {}).get('date', 'Unknown')}")
        c.drawString(50, height - 170, f"Probable Cause: {', '.join(incident_data.get('cause_keywords', []))}")
        
        # Section 2: Severity Assessment
        c.setFont("Helvetica-Bold", 14)
        c.drawString(50, height - 210, "⚠️ SEVERITY ASSESSMENT")
        c.setFont("Helvetica", 11)
        c.drawString(50, height - 230, f"Severity Score: {incident_data.get('severity_score', 0)}%")
        c.drawInlineImage(self.generate_severity_bar(incident_data.get('severity_score', 0)), 50, height - 250, width=400, height=20)
        
        # Section 3: Safety Recommendations
        c.setFont("Helvetica-Bold", 14)
        c.drawString(50, height - 290, "✅ SAFETY RECOMMENDATIONS")
        y = height - 310
        for tip in research_data.get('prevention_tips', []):
            c.drawString(70, y, f"• {tip}")
            y -= 15
            
        c.showPage()
        c.save()
        return file_path

    def generate_severity_bar(self, score: int) -> io.BytesIO:
        """Helper to create a severity bar image for PDF."""
        img = Image.new('RGB', (400, 20), color=(240, 240, 240))
        draw = ImageDraw.Draw(img)
        fill_color = (239, 68, 68) if score > 70 else (249, 115, 22) if score > 40 else (34, 197, 94)
        draw.rectangle([0, 0, (score/100)*400, 20], fill=fill_color)
        buf = io.BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)
        return buf

    def generate_poster(self, incident_data: dict, research_data: dict) -> str:
        """
        Generate an awareness poster with Pillow.
        """
        file_path = os.path.join(self.output_dir, f"Safety_Poster_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
        img = Image.new('RGB', (800, 1000), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        # Header background
        draw.rectangle([0, 0, 800, 150], fill=(59, 130, 246))
        
        # Text alignment
        title = "STOP ACCIDENTS. SAVE LIVES."
        draw.text((400, 75), title, fill=(255, 255, 255), anchor="mm", font_size=40)
        
        # Safety rules section
        draw.text((50, 200), "🚨 5 SAFETY RULES TO REMEMBER:", fill=(0, 0, 0), font_size=30)
        y = 250
        for i, tip in enumerate(research_data.get('prevention_tips', [])[:5]):
            draw.text((70, y), f"{i+1}. {tip}", fill=(50, 50, 50), font_size=20)
            y += 40
            
        # Statistics
        draw.rectangle([50, y+50, 750, y+150], outline=(239, 68, 68), width=2)
        draw.text((400, y+100), f"FACT: {research_data.get('statistics', {}).get('annual_deaths', '1.5L')} DEATHS / YEAR IN REGION", fill=(239, 68, 68), anchor="mm", font_size=25)
        
        img.save(file_path)
        return file_path

    def generate_document(self, incident_data: dict, research_data: dict) -> str:
        """
        Generate a Word document with python-docx.
        """
        file_path = os.path.join(self.output_dir, f"Safety_Guide_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.docx")
        doc = Document()
        doc.add_heading('COMPREHENSIVE ROAD SAFETY GUIDE', 0)
        
        doc.add_heading('1. Incident Overview', level=1)
        doc.add_paragraph(f"Based on Incident: {incident_data.get('location', {}).get('address', 'N/A')}")
        doc.add_paragraph(f"Date: {incident_data.get('datetime', {}).get('date', 'N/A')}")
        
        doc.add_heading('2. Key Statistics', level=1)
        table = doc.add_table(rows=1, cols=2)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Metric'
        hdr_cells[1].text = 'Value'
        
        stats = research_data.get('statistics', {})
        row_cells = table.add_row().cells
        row_cells[0].text = 'Annual Accidents'
        row_cells[1].text = str(stats.get('annual_accidents', 'N/A'))
        
        doc.add_heading('3. Prevention Tips', level=1)
        for tip in research_data.get('prevention_tips', []):
            doc.add_paragraph(tip, style='List Bullet')
            
        doc.save(file_path)
        return file_path

    def generate_presentation(self, incident_data: dict, research_data: dict) -> str:
        """
        Generate a PowerPoint presentation with python-pptx.
        """
        file_path = os.path.join(self.output_dir, f"Safety_Presentation_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
        prs = Presentation()
        
        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "ROAD SAFETY AWARENESS BRIEFING"
        subtitle.text = f"Incident at {incident_data.get('location', {}).get('address', 'N/A')}\nRoad Safety AI Agent"
        
        # Slide 2: Stats
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Regional Reality"
        tf = slide.placeholders[1].text_frame
        for cause in research_data.get('statistics', {}).get('common_causes', []):
            p = tf.add_paragraph()
            p.text = cause
            
        # Slide 3: Recommendations
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Prevention Strategies"
        tf = slide.placeholders[1].text_frame
        for tip in research_data.get('prevention_tips', []):
            p = tf.add_paragraph()
            p.text = tip
            
        prs.save(file_path)
        return file_path

    def generate_timeline(self, incident_data: dict) -> str:
        """
        Generate a horizontal timeline visual with matplotlib.
        """
        file_path = os.path.join(self.output_dir, f"Incident_Timeline_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
        
        # Simple data for timeline
        events = ["Before Incident", "Incident Moment", "Aftermath", "Response", "Reported"]
        dates = [0, 10, 20, 30, 40]
        
        plt.figure(figsize=(10, 2))
        plt.hlines(1, 0, 45, color='blue', alpha=0.5)
        plt.scatter(dates, [1]*len(dates), color='red', s=100, zorder=3)
        
        for i, event in enumerate(events):
            plt.text(dates[i], 1.05, event, rotation=45, ha='left')
            
        plt.axis('off')
        plt.tight_layout()
        plt.savefig(file_path)
        plt.close()
        return file_path
